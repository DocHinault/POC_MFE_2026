"""
Pipeline complet d'analyse des performances sociales
Récupère les KPI, envoie à GPT, génère PowerPoint et envoie par email
"""

import os
from datetime import datetime, timedelta
from typing import Dict, List, Any
import json
import logging
import requests
from io import BytesIO

logger = logging.getLogger(__name__)


class AnalysisPipeline:
    """Pipeline d'analyse complète des comptes sociaux"""
    
    def __init__(self, user_id: str, user_email: str, user_name: str):
        """
        Initialise le pipeline
        
        Args:
            user_id: ID du client
            user_email: Email du client
            user_name: Nom de l'entreprise
        """
        self.user_id = user_id
        self.user_email = user_email
        self.user_name = user_name
        self.kpis = {}
        self.gpt_recommendations = None
        
    def fetch_instagram_kpis(self, instagram_account_id: str, access_token: str) -> Dict[str, Any]:
        """
        Récupère les KPI Instagram du dernier mois via Instagram Graph API
        
        Args:
            instagram_account_id: ID du compte Instagram Business
            access_token: Token d'accès avec permissions
            
        Returns:
            Dict avec les KPI (impressions, reach, engagement, etc.)
        """
        try:
            base_url = "https://graph.instagram.com/v18.0"
            
            # Calcul de la période (dernier mois)
            end_date = datetime.now()
            start_date = end_date - timedelta(days=30)
            
            kpis = {
                'platform': 'Instagram',
                'account_id': instagram_account_id,
                'period_start': start_date.isoformat(),
                'period_end': end_date.isoformat(),
                'impressions': 0,
                'reach': 0,
                'engagement_rate': 0,
                'followers_growth': 0,
                'top_posts': [],
                'average_engagement': 0,
                'total_engagement': 0,
                'total_posts': 0,
                'profile_views': 0
            }
            
            # 1. Récupérer les insights du compte
            insights_fields = 'impressions,reach,profile_views,followers_count'
            insights_url = f"{base_url}/{instagram_account_id}/insights"
            params = {
                'metric': insights_fields,
                'period': 'day',
                'access_token': access_token
            }
            
            response = requests.get(insights_url, params=params)
            if response.status_code == 200:
                data = response.json().get('data', [])
                for item in data:
                    if item.get('name') == 'impressions':
                        kpis['impressions'] = item.get('total_value', {}).get('value', 0)
                    elif item.get('name') == 'reach':
                        kpis['reach'] = item.get('total_value', {}).get('value', 0)
                    elif item.get('name') == 'profile_views':
                        kpis['profile_views'] = item.get('total_value', {}).get('value', 0)
            
            # 2. Récupérer les posts du dernier mois
            media_url = f"{base_url}/{instagram_account_id}/media"
            media_params = {
                'fields': 'id,caption,media_type,media_product_type,timestamp,like_count,comments_count',
                'access_token': access_token,
                'limit': 100
            }
            
            response = requests.get(media_url, params=media_params)
            if response.status_code == 200:
                posts = response.json().get('data', [])
                
                # Filtrer les posts du dernier mois
                month_ago = datetime.now() - timedelta(days=30)
                recent_posts = []
                
                for post in posts:
                    post_date_str = post.get('timestamp', '')
                    if post_date_str:
                        post_date = datetime.fromisoformat(post_date_str.replace('Z', '+00:00'))
                        if post_date > month_ago:
                            engagement = post.get('like_count', 0) + post.get('comments_count', 0)
                            recent_posts.append({
                                'id': post.get('id'),
                                'caption': post.get('caption', '')[:100],
                                'likes': post.get('like_count', 0),
                                'comments': post.get('comments_count', 0),
                                'engagement': engagement,
                                'timestamp': post_date_str,
                                'type': post.get('media_type', '')
                            })
                            kpis['total_engagement'] += engagement
                
                kpis['total_posts'] = len(recent_posts)
                
                # Trier par engagement et prendre les top 5
                recent_posts.sort(key=lambda x: x['engagement'], reverse=True)
                kpis['top_posts'] = recent_posts[:5]
                
                # Calculer le taux d'engagement moyen
                if recent_posts:
                    kpis['average_engagement'] = kpis['total_engagement'] / len(recent_posts)
                
                # Calcul du taux d'engagement
                if kpis['reach'] > 0:
                    kpis['engagement_rate'] = (kpis['total_engagement'] / kpis['reach']) * 100
            
            logger.info(f"KPI Instagram récupérés: {kpis['total_posts']} posts, {kpis['reach']} reach")
            return kpis
            
        except Exception as e:
            logger.error(f"Erreur récupération Instagram KPI: {e}")
            return {}
    
    def fetch_facebook_kpis(self, page_id: str, access_token: str) -> Dict[str, Any]:
        """
        Récupère les KPI Facebook du dernier mois via Facebook Graph API
        
        Args:
            page_id: ID de la page Facebook
            access_token: Token d'accès avec permissions
            
        Returns:
            Dict avec les KPI (reach, engagement, fans growth, etc.)
        """
        try:
            base_url = "https://graph.facebook.com/v18.0"
            
            # Calcul de la période (dernier mois)
            end_date = datetime.now()
            start_date = end_date - timedelta(days=30)
            since = int(start_date.timestamp())
            until = int(end_date.timestamp())
            
            kpis = {
                'platform': 'Facebook',
                'page_id': page_id,
                'period_start': start_date.isoformat(),
                'period_end': end_date.isoformat(),
                'reach': 0,
                'impressions': 0,
                'engagement_rate': 0,
                'fans_growth': 0,
                'top_posts': [],
                'average_engagement': 0,
                'total_engagement': 0,
                'total_posts': 0,
                'page_views': 0,
                'total_fans': 0
            }
            
            # 1. Récupérer les insights de la page
            insights_url = f"{base_url}/{page_id}/insights"
            insights_metrics = 'page_post_engagements,page_impressions,page_fans,page_views_total'
            params = {
                'metric': insights_metrics,
                'period': 'day',
                'since': since,
                'until': until,
                'access_token': access_token
            }
            
            response = requests.get(insights_url, params=params)
            if response.status_code == 200:
                data = response.json().get('data', [])
                for item in data:
                    metric_name = item.get('name')
                    total_value = item.get('total_value', item.get('values', [{}])[0].get('value', 0))
                    
                    if metric_name == 'page_post_engagements':
                        kpis['total_engagement'] = total_value
                    elif metric_name == 'page_impressions':
                        kpis['impressions'] = total_value
                    elif metric_name == 'page_fans':
                        kpis['total_fans'] = total_value
                    elif metric_name == 'page_views_total':
                        kpis['page_views'] = total_value
            
            # 2. Récupérer les posts du dernier mois
            posts_url = f"{base_url}/{page_id}/posts"
            posts_params = {
                'fields': 'id,message,story,link,type,created_time,likes.summary(true),comments.summary(true),shares',
                'since': since,
                'until': until,
                'access_token': access_token,
                'limit': 100
            }
            
            response = requests.get(posts_url, params=posts_params)
            if response.status_code == 200:
                posts = response.json().get('data', [])
                
                top_posts = []
                for post in posts:
                    likes = post.get('likes', {}).get('summary', {}).get('total_count', 0)
                    comments = post.get('comments', {}).get('summary', {}).get('total_count', 0)
                    shares = post.get('shares', 0)
                    
                    engagement = likes + comments + (shares * 2)  # Les partages comptent pour 2
                    
                    top_posts.append({
                        'id': post.get('id'),
                        'message': (post.get('message') or post.get('story', ''))[:100],
                        'likes': likes,
                        'comments': comments,
                        'shares': shares,
                        'engagement': engagement,
                        'created_time': post.get('created_time'),
                        'type': post.get('type')
                    })
                
                kpis['total_posts'] = len(top_posts)
                
                # Trier par engagement et prendre les top 5
                top_posts.sort(key=lambda x: x['engagement'], reverse=True)
                kpis['top_posts'] = top_posts[:5]
                
                # Calculer le taux d'engagement moyen
                if top_posts:
                    kpis['average_engagement'] = kpis['total_engagement'] / len(top_posts)
                
                # Calcul du taux d'engagement
                if kpis['impressions'] > 0:
                    kpis['engagement_rate'] = (kpis['total_engagement'] / kpis['impressions']) * 100
            
            logger.info(f"KPI Facebook récupérés: {kpis['total_posts']} posts, {kpis['total_fans']} fans")
            return kpis
            
        except Exception as e:
            logger.error(f"Erreur récupération Facebook KPI: {e}")
            return {}
    
    def save_to_google_sheet(self, sheet_id: str, sheet_service) -> bool:
        """
        Sauvegarde les KPI dans la Google Sheet (onglet Analyse_Client)
        
        Args:
            sheet_id: ID de la Google Sheet
            sheet_service: Service Google Sheets (gspread ou google-api-python-client)
            
        Returns:
            True si succès, False sinon
        """
        try:
            from datetime import datetime
            
            # Créer les données à insérer
            row_data = [
                datetime.now().strftime('%Y-%m-%d %H:%M:%S'),  # Timestamp
                self.user_name,
                'Instagram',
                self.kpis.get('instagram', {}).get('impressions', 0),
                self.kpis.get('instagram', {}).get('reach', 0),
                self.kpis.get('instagram', {}).get('engagement_rate', 0),
                self.kpis.get('instagram', {}).get('total_posts', 0),
                self.kpis.get('instagram', {}).get('total_engagement', 0),
            ]
            
            facebook_data = [
                datetime.now().strftime('%Y-%m-%d %H:%M:%S'),  # Timestamp
                self.user_name,
                'Facebook',
                self.kpis.get('facebook', {}).get('impressions', 0),
                self.kpis.get('facebook', {}).get('reach', 0) if isinstance(self.kpis.get('facebook'), dict) else 0,
                self.kpis.get('facebook', {}).get('engagement_rate', 0),
                self.kpis.get('facebook', {}).get('total_posts', 0),
                self.kpis.get('facebook', {}).get('total_engagement', 0),
            ]
            
            # Optionnel: ajouter directement via gspread si disponible
            try:
                import gspread
                
                # Authentifier avec service account
                gc = gspread.service_account(filename='credentials.json')
                spreadsheet = gc.open_by_key(sheet_id)
                
                # Créer ou obtenir la feuille "Analyse_Client"
                try:
                    worksheet = spreadsheet.worksheet('Analyse_Client')
                except:
                    worksheet = spreadsheet.add_worksheet(title='Analyse_Client', rows=1000, cols=10)
                    # Ajouter les en-têtes
                    headers = ['Timestamp', 'Client', 'Platform', 'Impressions', 'Reach', 'Engagement %', 'Posts', 'Total Engagement']
                    worksheet.append_row(headers)
                
                # Ajouter les données
                worksheet.append_row(row_data)
                if 'facebook' in self.kpis:
                    worksheet.append_row(facebook_data)
                
                logger.info(f"✅ KPI sauvegardés dans Google Sheet {sheet_id}")
                return True
            except ImportError:
                logger.warning("gspread non disponible, skipping Google Sheet save")
                return False
            
        except Exception as e:
            logger.error(f"Erreur sauvegarde Google Sheet: {e}")
            return False
    
    def get_gpt_recommendations(self, gpt_api_key: str) -> Dict[str, Any]:
        """
        Envoie les KPI à GPT pour analysis et recommandations
        
        Args:
            gpt_api_key: Clé API OpenAI
            
        Returns:
            Dict avec:
            - objectives: Principaux objectifs
            - strengths: Points forts (+)
            - weaknesses: Points faibles (-)
            - next_post_ideas: 3 idées pour prochains posts
            - summary: Résumé général
        """
        try:
            from openai import OpenAI
            
            client = OpenAI(api_key=gpt_api_key)
            
            # Préparer le contexte KPI pour GPT
            kpi_context = json.dumps(self.kpis, indent=2, default=str)
            
            prompt = f"""
Tu es un expert en stratégie marketing et social media. Analyse les KPI suivants d'un compte social media pour le mois dernier et fournis des recommandations.

## KPI DU MOIS DERNIER:
{kpi_context}

## TÂCHE:
Analyse ces données et fournis une réponse JSON strictement structurée ainsi:

{{
    "objectives": ["Objectif 1", "Objectif 2", "Objectif 3"],
    "strengths": ["Force 1", "Force 2", "Force 3"],
    "weaknesses": ["Faiblesse 1", "Faiblesse 2", "Faiblesse 3"],
    "next_post_ideas": [
        {{"title": "Titre 1", "description": "Description courte", "expected_engagement": "high"}},
        {{"title": "Titre 2", "description": "Description courte", "expected_engagement": "high"}},
        {{"title": "Titre 3", "description": "Description courte", "expected_engagement": "medium"}}
    ],
    "summary": "Résumé général des performances"
}}

Réponds UNIQUEMENT avec le JSON, sans aucun texte supplémentaire.
"""
            
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.7,
                max_tokens=2000
            )
            
            # Parser la réponse JSON
            response_text = response.choices[0].message.content.strip()
            
            # Essayer de nettoyer la réponse (supprimer les marques markdown)
            if response_text.startswith('```'):
                response_text = response_text.split('```')[1]
                if response_text.startswith('json'):
                    response_text = response_text[4:]
            
            recommendations = json.loads(response_text)
            recommendations['timestamp'] = datetime.now().isoformat()
            recommendations['model'] = 'gpt-3.5-turbo'
            
            self.gpt_recommendations = recommendations
            logger.info("✅ Recommandations GPT générées")
            return recommendations
            
        except Exception as e:
            logger.error(f"Erreur appel GPT: {e}")
            # Retourner des recommandations par défaut
            return {
                'objectives': ['Augmenter l\'engagement', 'Croître la communauté'],
                'strengths': ['Contenu cohérent', 'Audience active'],
                'weaknesses': ['Améliorer la fréquence de publication'],
                'next_post_ideas': [
                    {'title': 'À générer', 'description': 'Erreur GPT'},
                    {'title': 'À générer', 'description': 'Erreur GPT'},
                    {'title': 'À générer', 'description': 'Erreur GPT'}
                ],
                'summary': f'Erreur lors de l\'analyse: {str(e)}'
            }
    
    def generate_powerpoint(self, output_path: str = None) -> str:
        """
        Génère un PowerPoint avec les insights et recommandations
        
        Args:
            output_path: Chemin de sauvegarde (optionnel)
            
        Returns:
            Chemin du fichier PowerPoint généré
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            # Définir un thème de couleurs
            title_color = RGBColor(31, 67, 153)  # Bleu foncé
            accent_color = RGBColor(66, 135, 245)  # Bleu clair
            text_color = RGBColor(50, 50, 50)  # Gris foncé
            
            def add_title_slide(prs, title, subtitle):
                """Ajoute une diapo de titre"""
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                
                # Background color
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(240, 245, 250)
                
                # Title
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
                title_frame = title_box.text_frame
                title_frame.word_wrap = True
                p = title_frame.paragraphs[0]
                p.text = title
                p.font.size = Pt(54)
                p.font.bold = True
                p.font.color.rgb = title_color
                
                # Subtitle
                subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(9), Inches(2))
                subtitle_frame = subtitle_box.text_frame
                p = subtitle_frame.paragraphs[0]
                p.text = subtitle
                p.font.size = Pt(28)
                p.font.color.rgb = accent_color
                
                # Date
                date_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.8))
                date_frame = date_box.text_frame
                p = date_frame.paragraphs[0]
                p.text = f"Rapport du {datetime.now().strftime('%d %B %Y')}"
                p.font.size = Pt(16)
                p.font.color.rgb = text_color
            
            def add_content_slide(prs, title, content_dict):
                """Ajoute une diapo de contenu"""
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                
                # Background
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 255)
                
                # Title
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
                title_frame = title_box.text_frame
                p = title_frame.paragraphs[0]
                p.text = title
                p.font.size = Pt(40)
                p.font.bold = True
                p.font.color.rgb = title_color
                
                # Content
                top_position = 1.3
                for key, value in content_dict.items():
                    box = slide.shapes.add_textbox(Inches(0.8), Inches(top_position), Inches(8.4), Inches(4))
                    text_frame = box.text_frame
                    text_frame.word_wrap = True
                    
                    # Key (bold)
                    p = text_frame.paragraphs[0]
                    p.text = f"• {key}"
                    p.font.size = Pt(18)
                    p.font.bold = True
                    p.font.color.rgb = accent_color
                    p.space_before = Pt(6)
                    
                    # Value
                    p = text_frame.add_paragraph()
                    p.text = str(value)
                    p.font.size = Pt(14)
                    p.font.color.rgb = text_color
                    p.level = 1
                    p.space_after = Pt(8)
                    
                    top_position += 0.8
            
            # Slide 1: Cover
            add_title_slide(
                prs,
                f"📊 Rapport d'Analyse",
                f"{self.user_name}"
            )
            
            # Slide 2: Executive Summary
            add_content_slide(prs, "📈 Résumé Exécutif", {
                'Période analysée': 'Dernier mois (30 jours)',
                'Plateformes': ', '.join([k.capitalize() for k in self.kpis.keys()]),
                'Résumé': self.gpt_recommendations.get('summary', 'Analyse complète réalisée')
            })
            
            # Slide 3: Instagram KPI
            if 'instagram' in self.kpis:
                ig = self.kpis['instagram']
                add_content_slide(prs, "📱 Instagram - Performances", {
                    'Impressions': f"{ig.get('impressions', 0):,}",
                    'Reach': f"{ig.get('reach', 0):,}",
                    'Total Engagement': f"{ig.get('total_engagement', 0):,}",
                    'Taux d\'engagement': f"{ig.get('engagement_rate', 0):.2f}%",
                    'Posts publiés': ig.get('total_posts', 0),
                    'Engagement moyen': f"{ig.get('average_engagement', 0):.1f}"
                })
            
            # Slide 4: Facebook KPI
            if 'facebook' in self.kpis:
                fb = self.kpis['facebook']
                add_content_slide(prs, "📘 Facebook - Performances", {
                    'Impressions': f"{fb.get('impressions', 0):,}",
                    'Total Fans': f"{fb.get('total_fans', 0):,}",
                    'Total Engagement': f"{fb.get('total_engagement', 0):,}",
                    'Taux d\'engagement': f"{fb.get('engagement_rate', 0):.2f}%",
                    'Posts publiés': fb.get('total_posts', 0),
                    'Engagement moyen': f"{fb.get('average_engagement', 0):.1f}"
                })
            
            # Slide 5: Strengths
            strengths_dict = {f"Atout {i+1}": s for i, s in enumerate(self.gpt_recommendations.get('strengths', [])[:4])}
            add_content_slide(prs, "💪 Points Forts", strengths_dict)
            
            # Slide 6: Weaknesses
            weaknesses_dict = {f"Point à améliorer {i+1}": w for i, w in enumerate(self.gpt_recommendations.get('weaknesses', [])[:4])}
            add_content_slide(prs, "📉 Points à Améliorer", weaknesses_dict)
            
            # Slide 7: Next Post Ideas
            ideas = self.gpt_recommendations.get('next_post_ideas', [])
            ideas_dict = {}
            for i, idea in enumerate(ideas[:3], 1):
                if isinstance(idea, dict):
                    ideas_dict[f"Idée {i}: {idea.get('title', 'Sans titre')}"] = idea.get('description', '')
                else:
                    ideas_dict[f"Idée {i}"] = str(idea)
            add_content_slide(prs, "💡 3 Idées pour les Prochains Posts", ideas_dict)
            
            # Slide 8: Recommendations
            objectives_dict = {f"Objectif {i+1}": o for i, o in enumerate(self.gpt_recommendations.get('objectives', [])[:4])}
            add_content_slide(prs, "🎯 Objectifs Recommandés", objectives_dict)
            
            # Save
            if output_path is None:
                output_path = f"/tmp/analysis_{self.user_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            
            prs.save(output_path)
            logger.info(f"✅ PowerPoint généré: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Erreur génération PowerPoint: {e}")
            return ""

    def generate_kpi_excel(self, output_path: str = None) -> str:
        """
        Génère un fichier Excel avec les KPI et les top posts.

        Args:
            output_path: Chemin de sortie (optionnel)

        Returns:
            Chemin du fichier Excel généré ou chaîne vide si échec
        """
        try:
            import pandas as pd

            if not self.kpis:
                logger.warning("⚠️ Aucun KPI disponible pour générer l'Excel")
                return ""

            if output_path is None:
                output_path = f"/tmp/kpis_{self.user_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"

            instagram_kpis = self.kpis.get('instagram') if isinstance(self.kpis.get('instagram'), dict) else None
            facebook_raw = self.kpis.get('facebook')
            facebook_kpis_list = []
            if isinstance(facebook_raw, list):
                facebook_kpis_list = facebook_raw
            elif isinstance(facebook_raw, dict) and facebook_raw:
                facebook_kpis_list = [facebook_raw]

            summary_rows = []
            if instagram_kpis:
                summary_rows.append({
                    'platform': 'Instagram',
                    'account_id': instagram_kpis.get('account_id'),
                    'period_start': instagram_kpis.get('period_start'),
                    'period_end': instagram_kpis.get('period_end'),
                    'impressions': instagram_kpis.get('impressions'),
                    'reach': instagram_kpis.get('reach'),
                    'engagement_rate': instagram_kpis.get('engagement_rate'),
                    'total_posts': instagram_kpis.get('total_posts'),
                    'total_engagement': instagram_kpis.get('total_engagement'),
                    'average_engagement': instagram_kpis.get('average_engagement'),
                    'profile_views': instagram_kpis.get('profile_views'),
                    'followers_growth': instagram_kpis.get('followers_growth')
                })

            for fb_kpis in facebook_kpis_list:
                summary_rows.append({
                    'platform': 'Facebook',
                    'page_id': fb_kpis.get('page_id'),
                    'period_start': fb_kpis.get('period_start'),
                    'period_end': fb_kpis.get('period_end'),
                    'impressions': fb_kpis.get('impressions'),
                    'reach': fb_kpis.get('reach'),
                    'engagement_rate': fb_kpis.get('engagement_rate'),
                    'total_posts': fb_kpis.get('total_posts'),
                    'total_engagement': fb_kpis.get('total_engagement'),
                    'average_engagement': fb_kpis.get('average_engagement'),
                    'page_views': fb_kpis.get('page_views'),
                    'total_fans': fb_kpis.get('total_fans'),
                    'fans_growth': fb_kpis.get('fans_growth')
                })

            with pd.ExcelWriter(output_path, engine="openpyxl") as writer:
                if summary_rows:
                    pd.DataFrame(summary_rows).to_excel(writer, sheet_name="KPI_Summary", index=False)

                if instagram_kpis and instagram_kpis.get('top_posts'):
                    pd.DataFrame(instagram_kpis.get('top_posts')).to_excel(
                        writer, sheet_name="Instagram_Top_Posts", index=False
                    )

                facebook_posts = []
                for fb_kpis in facebook_kpis_list:
                    for post in fb_kpis.get('top_posts', []):
                        facebook_posts.append({
                            'page_id': fb_kpis.get('page_id'),
                            **post
                        })
                if facebook_posts:
                    pd.DataFrame(facebook_posts).to_excel(
                        writer, sheet_name="Facebook_Top_Posts", index=False
                    )

            logger.info(f"✅ Excel KPI généré: {output_path}")
            return output_path

        except Exception as e:
            logger.error(f"Erreur génération Excel KPI: {e}")
            return ""
    
    def send_email_report(self, powerpoint_path: str, excel_path: str = None, sheet_url: str = None) -> bool:
        """
        Envoie le rapport par email interne avec PowerPoint et Excel KPI
        
        Args:
            powerpoint_path: Chemin du fichier PowerPoint
            excel_path: Chemin du fichier Excel KPI (optionnel)
            sheet_url: URL de la Google Sheet (optionnel)
            
        Returns:
            True si succès, False sinon
        """
        try:
            import smtplib
            from email.mime.text import MIMEText
            from email.mime.multipart import MIMEMultipart
            from email.mime.base import MIMEBase
            from email import encoders
            
            # Configuration email (à récupérer des variables d'environnement)
            sender_email = os.getenv('SMTP_EMAIL', 'noreply@example.com')
            sender_password = os.getenv('SMTP_PASSWORD', '')
            smtp_server = os.getenv('SMTP_SERVER', 'smtp.gmail.com')
            smtp_port = int(os.getenv('SMTP_PORT', '587'))
            recipient_email = os.getenv('ANALYSIS_REPORT_EMAIL')
            
            if not sender_password:
                logger.warning("❌ SMTP_PASSWORD non configuré, email non envoyé")
                return False
            if not recipient_email:
                logger.warning("❌ ANALYSIS_REPORT_EMAIL non configuré, email non envoyé")
                return False
            analysis = self.gpt_recommendations or {}
            
            # Créer le message
            msg = MIMEMultipart('alternative')
            msg['Subject'] = f"📊 Rapport d'Analyse - {self.user_name}"
            msg['From'] = sender_email
            msg['To'] = recipient_email
            
            # Corps du message HTML
            html_body = f"""
            <html>
                <body style="font-family: Arial, sans-serif; background-color: #f5f5f5; padding: 20px;">
                    <div style="background-color: white; padding: 30px; border-radius: 8px; max-width: 600px; margin: 0 auto;">
                        <h1 style="color: #1f4399; border-bottom: 3px solid #4287f5; padding-bottom: 10px;">
                            📊 Votre Rapport d'Analyse Social Media
                        </h1>
                        
                        <p style="color: #333; font-size: 16px;">
                            Bonjour <strong>{self.user_name}</strong>,
                        </p>
                        
                        <p style="color: #666; font-size: 14px; line-height: 1.6;">
                            Nous avons analysé vos performances sur les réseaux sociaux du mois dernier.
                            Vous trouverez ci-joint votre rapport détaillé avec recommandations.
                        </p>
                        
                        <hr style="border: none; border-top: 2px solid #e0e0e0; margin: 20px 0;">
                        
                        <h2 style="color: #1f4399; margin-top: 25px;">🎯 Principaux Objectifs</h2>
                        <ul style="color: #333; font-size: 14px;">
            """
            
            for obj in analysis.get('objectives', []):
                html_body += f"<li style='margin: 8px 0;'>{obj}</li>"
            
            html_body += """
                        </ul>
                        
                        <h2 style="color: #1f4399;">💪 Vos Points Forts</h2>
                        <ul style="color: #333; font-size: 14px;">
            """
            
            for strength in analysis.get('strengths', []):
                html_body += f"<li style='margin: 8px 0;'>{strength}</li>"
            
            html_body += """
                        </ul>
                        
                        <h2 style="color: #1f4399;">📉 Points à Améliorer</h2>
                        <ul style="color: #333; font-size: 14px;">
            """
            
            for weakness in analysis.get('weaknesses', []):
                html_body += f"<li style='margin: 8px 0;'>{weakness}</li>"
            
            html_body += """
                        </ul>
                        
                        <h2 style="color: #1f4399;">💡 3 Idées pour vos Prochains Posts</h2>
                        <ol style="color: #333; font-size: 14px;">
            """
            
            for idea in analysis.get('next_post_ideas', []):
                if isinstance(idea, dict):
                    idea_text = f"<strong>{idea.get('title', 'Idée')}</strong>: {idea.get('description', '')}"
                else:
                    idea_text = str(idea)
                html_body += f"<li style='margin: 8px 0;'>{idea_text}</li>"
            
            html_body += f"""
                        </ol>
                        
                        <hr style="border: none; border-top: 2px solid #e0e0e0; margin: 20px 0;">
                        
                        <p style="color: #666; font-size: 12px; text-align: center; margin-top: 30px;">
                            📎 <strong>Pièces jointes:</strong> Présentation PowerPoint + Excel KPI<br>
                            Généré le {datetime.now().strftime('%d/%m/%Y à %H:%M')}
                        </p>
                        
                        <p style="color: #999; font-size: 11px; text-align: center;">
                            {"" if not sheet_url else f"<a href='{sheet_url}' style='color: #4287f5; text-decoration: none;'>Accéder à la Google Sheet</a><br>"}
                            © 2026 - MG Social Media Dashboard
                        </p>
                    </div>
                </body>
            </html>
            """
            
            # Attacher le corps HTML
            msg.attach(MIMEText(html_body, 'html'))
            
            # Attacher le PowerPoint
            if powerpoint_path and os.path.exists(powerpoint_path):
                with open(powerpoint_path, 'rb') as attachment:
                    part = MIMEBase('application', 'octet-stream')
                    part.set_payload(attachment.read())
                    encoders.encode_base64(part)
                    part.add_header(
                        'Content-Disposition',
                        f'attachment; filename= {os.path.basename(powerpoint_path)}'
                    )
                    msg.attach(part)

            # Attacher l'Excel KPI
            if excel_path and os.path.exists(excel_path):
                with open(excel_path, 'rb') as attachment:
                    part = MIMEBase('application', 'octet-stream')
                    part.set_payload(attachment.read())
                    encoders.encode_base64(part)
                    part.add_header(
                        'Content-Disposition',
                        f'attachment; filename= {os.path.basename(excel_path)}'
                    )
                    msg.attach(part)
            
            # Envoyer l'email
            server = smtplib.SMTP(smtp_server, smtp_port)
            server.starttls()
            server.login(sender_email, sender_password)
            server.send_message(msg)
            server.quit()
            
            logger.info(f"✅ Email envoyé à {recipient_email}")
            return True
            
        except Exception as e:
            logger.error(f"Erreur envoi email: {e}")
            return False
    
    def run_full_pipeline(self, instagram_data: Dict = None, facebook_data: List[Dict] = None, 
                         sheet_id: str = None) -> Dict[str, Any]:
        """
        Lance le pipeline complet d'analyse
        
        Args:
            instagram_data: Dict avec id, access_token du compte Instagram
            facebook_data: Liste de dicts avec id, access_token des pages Facebook
            sheet_id: ID de la Google Sheet (optionnel)
            
        Returns:
            Dict avec le résultat de chaque étape et les chemins des fichiers générés
        """
        result = {
            'success': False,
            'instagram_kpis': None,
            'facebook_kpis': None,
            'gpt_recommendations': None,
            'powerpoint_path': None,
            'excel_path': None,
            'sheet_saved': False,
            'email_sent': False,
            'errors': []
        }
        
        try:
            logger.info(f"🚀 Démarrage du pipeline pour {self.user_name}")
            
            # Étape 1: Récupérer les KPI Instagram
            if instagram_data and instagram_data.get('id') and instagram_data.get('access_token'):
                logger.info("📸 Récupération des KPI Instagram...")
                try:
                    self.kpis['instagram'] = self.fetch_instagram_kpis(
                        instagram_data['id'],
                        instagram_data['access_token']
                    )
                    result['instagram_kpis'] = self.kpis.get('instagram')
                    logger.info(f"✅ Instagram KPI récupérés: {self.kpis['instagram'].get('total_posts', 0)} posts")
                except Exception as e:
                    result['errors'].append(f"Instagram: {str(e)}")
                    logger.error(f"❌ Erreur Instagram: {e}")
            
            # Étape 2: Récupérer les KPI Facebook
            if facebook_data:
                logger.info("📘 Récupération des KPI Facebook...")
                try:
                    self.kpis['facebook'] = []
                    for page in facebook_data:
                        if page.get('id') and page.get('access_token'):
                            fb_kpi = self.fetch_facebook_kpis(
                                page['id'],
                                page['access_token']
                            )
                            self.kpis['facebook'].append(fb_kpi)
                    
                    if self.kpis['facebook']:
                        result['facebook_kpis'] = self.kpis['facebook']
                        logger.info(f"✅ Facebook KPI récupérés: {len(self.kpis['facebook'])} page(s)")
                except Exception as e:
                    result['errors'].append(f"Facebook: {str(e)}")
                    logger.error(f"❌ Erreur Facebook: {e}")
            
            # Arrêter si aucun KPI ne peut être récupéré
            if not self.kpis or (not self.kpis.get('instagram') and not self.kpis.get('facebook')):
                result['errors'].append("Aucun KPI disponible à analyser")
                logger.error("❌ Aucun KPI disponible")
                return result
            
            # Étape 3: Sauvegarder dans Google Sheet
            if sheet_id:
                logger.info("💾 Sauvegarde dans Google Sheet...")
                try:
                    if self.save_to_google_sheet(sheet_id, None):
                        result['sheet_saved'] = True
                        logger.info("✅ Google Sheet mise à jour")
                except Exception as e:
                    result['errors'].append(f"Google Sheet: {str(e)}")
                    logger.warning(f"⚠️ Google Sheet non sauvegardée: {e}")
            
            # Étape 4: Obtenir recommandations GPT
            gpt_key = os.getenv('OPENAI_API_KEY')
            if gpt_key:
                logger.info("🤖 Analyse GPT...")
                try:
                    recs = self.get_gpt_recommendations(gpt_key)
                    result['gpt_recommendations'] = recs
                    logger.info("✅ Recommandations GPT générées")
                except Exception as e:
                    result['errors'].append(f"GPT: {str(e)}")
                    logger.warning(f"⚠️ GPT erreur: {e}")
                    # Continuer avec des recommandations par défaut
            
            # Étape 5: Générer PowerPoint
            if self.gpt_recommendations or self.kpis:
                logger.info("📊 Génération PowerPoint...")
                try:
                    pptx_path = self.generate_powerpoint()
                    if pptx_path:
                        result['powerpoint_path'] = pptx_path
                        logger.info(f"✅ PowerPoint généré: {pptx_path}")
                except Exception as e:
                    result['errors'].append(f"PowerPoint: {str(e)}")
                    logger.error(f"❌ PowerPoint erreur: {e}")

            # Étape 5bis: Générer Excel KPI
            if self.kpis:
                logger.info("📈 Génération Excel KPI...")
                try:
                    excel_path = self.generate_kpi_excel()
                    if excel_path:
                        result['excel_path'] = excel_path
                        logger.info(f"✅ Excel KPI généré: {excel_path}")
                except Exception as e:
                    result['errors'].append(f"Excel KPI: {str(e)}")
                    logger.error(f"❌ Excel KPI erreur: {e}")
            
            # Étape 6: Envoyer par email
            if result.get('powerpoint_path') and result.get('excel_path'):
                logger.info("📧 Envoi du rapport par email...")
                try:
                    email_sent = self.send_email_report(
                        result['powerpoint_path'],
                        excel_path=result.get('excel_path'),
                        sheet_url=None
                    )
                    result['email_sent'] = email_sent
                    if email_sent:
                        logger.info(f"✅ Email envoyé à {self.user_email}")
                except Exception as e:
                    result['errors'].append(f"Email: {str(e)}")
                    logger.error(f"❌ Email erreur: {e}")
            
            result['success'] = len(result['errors']) == 0 or result.get('powerpoint_path') is not None
            
            if result['success']:
                logger.info("✅ Pipeline complété avec succès!")
            else:
                logger.warning(f"⚠️ Pipeline complété avec erreurs: {result['errors']}")
            
            return result
            
        except Exception as e:
            result['errors'].append(f"Pipeline: {str(e)}")
            logger.error(f"Erreur pipeline global: {e}")
            return result
